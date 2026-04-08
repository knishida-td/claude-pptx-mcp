#!/usr/bin/env python3
"""
SlideKit diagram helpers for python-pptx.

pptx_generate で作成したPPTXにネイティブ図解（円・矩形・矢印等）を追加するための
ヘルパー関数とバリデーション関数。

Usage:
    from diagram_helpers import DiagramKit
    kit = DiagramKit('/tmp/input.pptx')
    sl = kit.slide(5)  # 0-indexed ではなくスライド番号
    kit.rrect(sl, 0.5, 1.0, 4.0, 3.0, fc=kit.ORANGE, lc=kit.YELLOW, lw=2)
    kit.save('/tmp/output.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class DiagramKit:
    """SlideKit 図解ヘルパー"""

    # === SlideKit カラーパレット ===
    ORANGE = RGBColor(0xEF, 0x48, 0x23)
    YELLOW = RGBColor(0xFC, 0xBF, 0x17)
    DARK = RGBColor(0x33, 0x33, 0x33)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x88, 0x88, 0x88)
    LGRAY = RGBColor(0xCC, 0xCC, 0xCC)
    TEAL = RGBColor(0x4E, 0xCD, 0xC4)
    LIGHT_ORANGE = RGBColor(0xFF, 0xF5, 0xF0)
    RED_LIGHT = RGBColor(0xFF, 0x6B, 0x6B)

    # === フェーズ背景色 ===
    PHASE_INPUT = RGBColor(0xFF, 0xE8, 0xE0)
    PHASE_CONVERT = RGBColor(0xFF, 0xF3, 0xD0)
    PHASE_OUTPUT = RGBColor(0xE0, 0xF5, 0xF3)

    # === SlideKit レイアウト定数 ===
    CONTENT_TOP = 0.893      # コンテンツ領域の上端 (inch)
    CONTENT_BOTTOM = 4.837   # コンテンツ領域の下端 (inch)
    CONTENT_HEIGHT = 3.944   # コンテンツ領域の高さ (inch)
    SLIDE_WIDTH = 10.0       # スライド幅 (inch)

    FONT = 'Hiragino Kaku Gothic Pro W3'

    # テンプレートシェイプ名（削除禁止）
    TEMPLATE_SHAPES = {'Text 0', 'Shape 1', 'Shape 2', 'Text 7', 'Text 8'}

    def __init__(self, path: str):
        self.prs = Presentation(path)

    def slide(self, num: int):
        """スライド番号（1始まり）でスライドを取得"""
        return self.prs.slides[num - 1]

    def save(self, path: str):
        self.prs.save(path)
        print(f"Saved to {path}")

    # ------------------------------------------------------------------
    # 縦中央配置の計算
    # ------------------------------------------------------------------
    def ideal_top(self, content_height: float) -> float:
        """コンテンツの理想的なy座標を計算（inch）"""
        return self.CONTENT_TOP + (self.CONTENT_HEIGHT - content_height) / 2

    # ------------------------------------------------------------------
    # テンプレートシェイプの管理
    # ------------------------------------------------------------------
    def clear_custom_shapes(self, sl):
        """テンプレートシェイプ以外を全削除"""
        to_remove = [s for s in sl.shapes if s.name not in self.TEMPLATE_SHAPES]
        sp_tree = sl.shapes._spTree
        for s in to_remove:
            sp_tree.remove(s._element)
        return len(to_remove)

    def clear_bigtext_placeholders(self, sl):
        """bigtextレイアウトの空プレースホルダーのみ削除
        (Text 3: 空ヘッダー, Shape 4: 区切り線, Text 5: 空サブテキスト)"""
        removed = 0
        sp_tree = sl.shapes._spTree
        for s in list(sl.shapes):
            if s.name == 'Text 3' and abs(s.width / 914400 - 9.0) < 0.5:
                if not s.text_frame.text.strip():
                    sp_tree.remove(s._element); removed += 1
            elif s.name == 'Shape 4' and s.height / 914400 < 0.02:
                sp_tree.remove(s._element); removed += 1
            elif s.name == 'Text 5' and abs(s.width / 914400 - 9.0) < 0.5:
                if not s.text_frame.text.strip():
                    sp_tree.remove(s._element); removed += 1
        return removed

    # ------------------------------------------------------------------
    # 基本シェイプ
    # ------------------------------------------------------------------
    def rrect(self, sl, x, y, w, h, fc=None, lc=None, lw=0, adj=0.05):
        """角丸矩形を追加。座標はinch。"""
        fc = fc or self.WHITE
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        if lc:
            s.line.color.rgb = lc; s.line.width = Pt(lw)
        else:
            s.line.fill.background()
        s.adjustments[0] = adj
        return s

    def rect(self, sl, x, y, w, h, fc=None):
        """矩形を追加。座標はinch。"""
        fc = fc or self.WHITE
        s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        s.line.fill.background()
        return s

    def circ(self, sl, x, y, d, fc):
        """円を追加。座標・直径はinch。"""
        s = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(x), Inches(y), Inches(d), Inches(d))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        s.line.fill.background()
        return s

    def rarr(self, sl, x, y, w, h, fc=None):
        """右矢印。座標はinch。"""
        fc = fc or self.DARK
        s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        s.line.fill.background()
        return s

    def darr(self, sl, x, y, w, h, fc=None):
        """下矢印。座標はinch。"""
        fc = fc or self.DARK
        s = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        s.line.fill.background()
        return s

    def larr(self, sl, x, y, w, h, fc=None):
        """左矢印。座標はinch。"""
        fc = fc or self.DARK
        s = sl.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                                Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fc
        s.line.fill.background()
        return s

    # ------------------------------------------------------------------
    # テキスト操作
    # ------------------------------------------------------------------
    def stxt(self, shape, txt, fs=14, c=None, b=False, al=PP_ALIGN.CENTER):
        """シェイプ内にテキストを設定。fsはpt値。"""
        c = c or self.DARK
        tf = shape.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = al
        p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run()
        r.text = txt; r.font.size = Pt(fs)
        r.font.color.rgb = c; r.font.bold = b; r.font.name = self.FONT
        return tf

    def stxt2(self, shape, t1, fs1, c1, b1, t2, fs2, c2, b2=False):
        """シェイプ内に2行テキスト。"""
        c1 = c1 or self.DARK; c2 = c2 or self.DARK
        tf = shape.text_frame; tf.clear(); tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = t1; r.font.size = Pt(fs1)
        r.font.color.rgb = c1; r.font.bold = b1; r.font.name = self.FONT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(0)
        r2 = p2.add_run()
        r2.text = t2; r2.font.size = Pt(fs2)
        r2.font.color.rgb = c2; r2.font.bold = b2; r2.font.name = self.FONT
        return tf

    def tbox(self, sl, x, y, w, h, txt, fs=14, c=None, b=False,
             al=PP_ALIGN.CENTER):
        """テキストボックスを追加。座標はinch。"""
        c = c or self.DARK
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = al
        p.space_before = Pt(0); p.space_after = Pt(0)
        r = p.add_run()
        r.text = txt; r.font.size = Pt(fs)
        r.font.color.rgb = c; r.font.bold = b; r.font.name = self.FONT
        return tb

    # ------------------------------------------------------------------
    # 背面配置
    # ------------------------------------------------------------------
    def send_to_back(self, sl, shape):
        """シェイプをスライドの最背面に移動"""
        sp_tree = sl.shapes._spTree
        elem = shape._element
        sp_tree.remove(elem)
        sp_tree.insert(2, elem)  # nvGrpSpPr, grpSpPr の後

    # ------------------------------------------------------------------
    # KeyMsg背景の確認・追加
    # ------------------------------------------------------------------
    def ensure_keymsg_bg(self, sl):
        """KeyMsg背景(FFF5F0)がなければ追加"""
        keymsg_text = None
        has_bg = False
        for s in sl.shapes:
            y_in = s.top / 914400
            w_in = s.width / 914400
            h_in = s.height / 914400
            if y_in > 4.5 and w_in > 8.0 and h_in < 0.6:
                try:
                    if str(s.fill.fore_color.rgb) == 'FFF5F0':
                        has_bg = True; continue
                except:
                    pass
                if s.has_text_frame and s.text_frame.text.strip():
                    keymsg_text = s

        if keymsg_text and not has_bg:
            bg = sl.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(4.837), Inches(9.0), Inches(0.4))
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.LIGHT_ORANGE
            bg.line.fill.background()
            # Text 7 の前に挿入
            sp_tree = sl.shapes._spTree
            bg_elem = bg._element
            txt_elem = keymsg_text._element
            sp_tree.remove(bg_elem)
            sp_tree.insert(list(sp_tree).index(txt_elem), bg_elem)
            return True
        return False

    # ------------------------------------------------------------------
    # バリデーション
    # ------------------------------------------------------------------
    def check_overflow(self, sl, slide_num=None):
        """スライド内の全テキストの改行オーバーフローをチェック"""
        issues = []
        for s in sl.shapes:
            if not s.has_text_frame:
                continue
            shape_w_pt = s.width / 12700
            usable = shape_w_pt - 14.4  # 内部マージン

            for p in s.text_frame.paragraphs:
                full = ''.join(r.text for r in p.runs)
                if not full.strip():
                    continue
                fs = None
                for r in p.runs:
                    if r.font.size:
                        fs = r.font.size / 12700
                        break
                if not fs:
                    continue

                for line in full.split('\n'):
                    if not line.strip() or len(line) <= 3:
                        continue
                    est_w = self._estimate_width(line, fs)
                    if est_w > usable:
                        ovf = (est_w - usable) / usable * 100
                        if ovf > 3:
                            issues.append({
                                'slide': slide_num,
                                'shape': s.name,
                                'text': line[:50],
                                'font_pt': fs,
                                'box_w_in': s.width / 914400,
                                'overflow_pct': ovf,
                            })
        return issues

    def check_all_overflow(self):
        """全スライドの改行オーバーフローをチェック"""
        all_issues = []
        for i, sl in enumerate(self.prs.slides):
            all_issues.extend(self.check_overflow(sl, i + 1))
        return all_issues

    def check_bounds(self):
        """全スライドのシェイプがスライド境界外に出ていないかチェック"""
        issues = []
        for i, sl in enumerate(self.prs.slides):
            for s in sl.shapes:
                x = s.left / 914400
                y = s.top / 914400
                w = s.width / 914400
                h = s.height / 914400
                if x + w > 10.1 or y + h > 5.7 or x < -0.1 or y < -0.1:
                    issues.append({
                        'slide': i + 1,
                        'shape': s.name,
                        'x': x, 'y': y, 'w': w, 'h': h,
                        'right': x + w, 'bottom': y + h,
                    })
        return issues

    @staticmethod
    def _estimate_width(text: str, fs_pt: float) -> float:
        """テキスト幅をpt単位で推定"""
        w = 0
        for ch in text:
            if ord(ch) > 0x2FFF or ch in '「」（）、。・：％→＝×':
                w += fs_pt * 1.05
            elif ch in '0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz':
                w += fs_pt * 0.6
            elif ch == ' ':
                w += fs_pt * 0.3
            else:
                w += fs_pt * 0.7
        return w


# ======================================================================
# 図解パターン（コピペ用テンプレート）
# ======================================================================

def pattern_horizontal_flow(kit, sl, y0, steps, arrow_color=None):
    """横フロー（3ステップまで）
    steps: [(label, subtitle, color, text_color), ...]
    """
    arrow_color = arrow_color or kit.LGRAY
    n = len(steps)
    step_w = 2.5
    arrow_w = 0.5
    total = n * step_w + (n - 1) * (arrow_w + 0.2)
    sx = (10 - total) / 2

    for i, (label, subtitle, color, tc) in enumerate(steps):
        x = sx + i * (step_w + arrow_w + 0.2)
        box = kit.rrect(sl, x, y0, step_w, 1.2, fc=color, adj=0.08)
        kit.stxt2(box, label, 18, tc, True, subtitle, 11, tc)
        if i < n - 1:
            ax = x + step_w + 0.1
            kit.rarr(sl, ax, y0 + 0.4, arrow_w, 0.25, fc=arrow_color)


def pattern_card_grid(kit, sl, y0, cards, cols=3, card_h=1.5,
                      border_color=None):
    """カードグリッド（2×3, 1×4 等）
    cards: [(title, description), ...]
    """
    border_color = border_color or kit.LGRAY
    gap = 0.15
    card_w = (9.0 - gap * (cols - 1)) / cols
    sx = 0.5

    for i, (title, desc) in enumerate(cards):
        col, row = i % cols, i // cols
        x = sx + col * (card_w + gap)
        y = y0 + row * (card_h + gap)
        c = kit.rrect(sl, x, y, card_w, card_h, lc=border_color, lw=1.5)
        kit.stxt2(c, title, 13, kit.ORANGE, True, desc, 11, kit.GRAY)


def pattern_two_column(kit, sl, y0, left_title, left_items,
                       right_title, right_items,
                       left_color=None, right_color=None):
    """2カラム比較ボックス
    items: [(heading, description), ...]
    """
    left_color = left_color or kit.ORANGE
    right_color = right_color or kit.YELLOW
    bw, bh = 4.1, 3.0

    lb = kit.rrect(sl, 0.5, y0, bw, bh, lc=left_color, lw=2)
    rb = kit.rrect(sl, 5.0, y0, bw, bh, lc=right_color, lw=2)

    # 丸アイコン
    cd = 0.55
    lcx = 0.5 + (bw - cd) / 2
    rcx = 5.0 + (bw - cd) / 2
    lc = kit.circ(sl, lcx, y0 + 0.1, cd, left_color)
    ltc = kit.WHITE if left_color == kit.ORANGE else kit.DARK
    kit.stxt(lc, left_title, 11, ltc, True)
    rc = kit.circ(sl, rcx, y0 + 0.1, cd, right_color)
    rtc = kit.WHITE if right_color != kit.YELLOW else kit.DARK
    kit.stxt(rc, right_title, 11, rtc, True)

    # 箇条書き
    for i, (h, d) in enumerate(left_items):
        iy = y0 + 0.8 + i * 0.7
        kit.tbox(sl, 0.9, iy, 3.3, 0.3, f'● {h}', 13, kit.DARK, True, PP_ALIGN.LEFT)
        kit.tbox(sl, 1.15, iy + 0.28, 3.3, 0.25, d, 11, kit.GRAY, False, PP_ALIGN.LEFT)
    for i, (h, d) in enumerate(right_items):
        iy = y0 + 0.8 + i * 0.7
        kit.tbox(sl, 5.4, iy, 3.3, 0.3, f'● {h}', 13, kit.DARK, True, PP_ALIGN.LEFT)
        kit.tbox(sl, 5.65, iy + 0.28, 3.3, 0.25, d, 11, kit.GRAY, False, PP_ALIGN.LEFT)


def pattern_chip_grid(kit, sl, y0, categories, chips_per_row=6):
    """チップ一覧（カテゴリ×N語）
    categories: [(name, color, text_color, [words...]), ...]
    """
    label_w = 0.7
    h_gap = 0.10
    v_gap = 0.10
    chip_h = 0.40
    chips_start = 0.5 + label_w + h_gap
    avail = 9.0 - label_w - h_gap
    chip_w = (avail - h_gap * (chips_per_row - 1)) / chips_per_row

    for row, (name, color, tc, words) in enumerate(categories):
        ry = y0 + row * (chip_h + v_gap)
        lbl = kit.rrect(sl, 0.5, ry, label_w, chip_h, fc=color, adj=0.2)
        kit.stxt(lbl, name, 9, tc, True)
        for col, word in enumerate(words[:chips_per_row]):
            cx = chips_start + col * (chip_w + h_gap)
            chip = kit.rrect(sl, cx, ry, chip_w, chip_h, lc=color, lw=1.2, adj=0.2)
            kit.stxt(chip, word, 10, kit.DARK, False)


def pattern_person_cards(kit, sl, y0, fact_text, persons):
    """ファクトバー + 人物カード
    persons: [(role, question, message, color, text_color), ...]
    """
    n = len(persons)
    card_w = (9.0 - 0.2 * (n - 1)) / n
    sx = 0.5

    # ファクトバー
    fact = kit.rrect(sl, 0.5, y0, 9.0, 0.5, fc=kit.DARK, adj=0.1)
    kit.stxt(fact, fact_text, 13, kit.WHITE, True)

    card_y = y0 + 0.7
    for i, (role, question, message, color, tc) in enumerate(persons):
        x = sx + i * (card_w + 0.2)
        card = kit.rrect(sl, x, card_y, card_w, 2.7, lc=color, lw=2)
        # 下矢印
        kit.darr(sl, x + card_w / 2 - 0.2, y0 + 0.5, 0.4, 0.2, fc=kit.LGRAY)
        # 人物アイコン
        cd = 0.6
        pc = kit.circ(sl, x + (card_w - cd) / 2, card_y + 0.15, cd, color)
        kit.stxt(pc, role, 11, tc, True)
        # 質問
        kit.tbox(sl, x, card_y + 0.85, card_w, 0.5, question, 11, kit.GRAY)
        # メッセージ
        mb = kit.rrect(sl, x + 0.15, card_y + 1.45, card_w - 0.3, 1.1,
                       fc=RGBColor(0xF9, 0xF9, 0xF9), lc=color, lw=1)
        kit.stxt(mb, message, 11, kit.DARK, False)
